"""Build the final solution presentation for the LG Aimers Phase 2 entry."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "output" / "LG_Aimers_솔루션_PPT_Phase2.pptx"
SOLUTION_PPTX = ROOT / "solution" / "LG_Aimers_솔루션_PPT_Phase2.pptx"

W = Inches(13.333333)
H = Inches(7.5)

BG = "0B1020"
PANEL = "151C31"
PANEL_2 = "1C2640"
INK = "F5F7FB"
MUTED = "AEB8CC"
TEAL = "39D8C6"
TEAL_DARK = "166E6B"
CORAL = "FF6B6B"
GOLD = "F4C95D"
BLUE = "6EA8FE"
GRID = "2A3655"

FONT = "Noto Sans CJK KR"
FONT_MONO = "Noto Sans Mono CJK KR"


def rgb(hex_code: str) -> RGBColor:
    return RGBColor.from_string(hex_code)


def add_rect(slide, x, y, w, h, fill, radius=False, line=None, width=1):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(width)
    else:
        shape.line.fill.background()
    if radius:
        try:
            shape.adjustments[0] = 0.12
        except (IndexError, ValueError):
            pass
    return shape


def add_line(slide, x1, y1, x2, y2, color=GRID, width=1.2, dash=None):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dash is not None:
        line.line.dash_style = dash
    return line


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=20,
    color=INK,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.02,
    line_spacing=1.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.line_spacing = line_spacing
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_label(slide, text, x, y, w, color=TEAL, fill=PANEL_2):
    add_rect(slide, x, y, w, 0.32, fill, radius=True)
    add_text(slide, text, x, y + 0.01, w, 0.27, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)


def add_bullet_list(slide, items, x, y, w, h, size=16, color=INK, bullet_color=TEAL, gap=0.43):
    for idx, item in enumerate(items):
        yy = y + idx * gap
        add_text(slide, "•", x, yy - 0.01, 0.25, 0.3, size=size + 2, color=bullet_color, bold=True)
        add_text(slide, item, x + 0.27, yy, w - 0.27, min(gap, h), size=size, color=color)


def add_header(slide, number, eyebrow, title, subtitle=None):
    add_text(slide, f"{number:02d}", 0.55, 0.40, 0.45, 0.28, size=11, color=TEAL, bold=True)
    add_text(slide, eyebrow.upper(), 1.05, 0.40, 4.5, 0.28, size=10, color=MUTED, bold=True)
    add_text(slide, title, 0.55, 0.80, 12.2, 0.60, size=26, color=INK, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.58, 1.43, 12.0, 0.40, size=12, color=MUTED)
    add_line(slide, 0.55, 1.92, 12.78, 1.92, GRID, 0.8)


def add_footer(slide, number):
    add_text(slide, "LG Aimers 9기 · 제구 성공 확률 예측 · LG Aimers 9기", 0.55, 7.15, 6.0, 0.20, size=8, color="74809A")
    add_text(slide, f"{number:02d}", 12.30, 7.15, 0.45, 0.20, size=8, color="74809A", align=PP_ALIGN.RIGHT)


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(BG)


def add_card_title(slide, title, x, y, w, color=TEAL):
    add_text(slide, title, x, y, w, 0.34, size=13, color=color, bold=True)


def build_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # 01 — Cover
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_rect(slide, 8.50, -0.15, 5.10, 7.80, "10182C")
    add_rect(slide, 8.92, 0.46, 3.65, 3.65, BG, line=GRID, width=1.4)
    diamond = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DIAMOND, Inches(9.55), Inches(1.06), Inches(2.4), Inches(2.4))
    diamond.fill.background()
    diamond.line.color.rgb = rgb(TEAL_DARK)
    diamond.line.width = Pt(2.3)
    for bx, by in [(10.67, 0.92), (11.83, 2.06), (10.67, 3.20), (9.52, 2.06)]:
        b = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(bx), Inches(by), Inches(0.26), Inches(0.26))
        b.fill.solid(); b.fill.fore_color.rgb = rgb(TEAL); b.line.fill.background()
    ball = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(10.30), Inches(4.64), Inches(1.05), Inches(1.05))
    ball.fill.solid(); ball.fill.fore_color.rgb = rgb(INK); ball.line.color.rgb = rgb("C8D1E2")
    add_line(slide, 10.55, 4.78, 11.06, 5.55, CORAL, 1.4)
    add_line(slide, 11.12, 4.80, 10.63, 5.54, CORAL, 1.4)
    add_label(slide, "LG AIMERS 9기 · PHASE 2 FINAL", 0.63, 0.57, 2.75)
    add_text(slide, "투구 직전 정보로\n제구 성공 확률을 읽다", 0.63, 1.35, 7.45, 1.75, size=36, bold=True)
    add_text(slide, "3-Tier Multi-Family GBDT Super Ensemble & Adaptive Hierarchical Gate\n최종 제출: submit_ref4_super113A.zip", 0.68, 3.38, 7.55, 0.78, size=15, color=MUTED)
    add_rect(slide, 0.65, 4.45, 7.10, 0.04, TEAL)
    add_text(slide, "LG AIMERS 9기 PHASE 2", 0.68, 4.85, 6.6, 0.36, size=13, color=INK, bold=True)
    add_text(slide, "공식 최종 점수: 1121.9039933605점  |  Public 180위", 0.68, 5.30, 7.1, 0.34, size=14, color=GOLD, bold=True)
    add_text(slide, "학습 데이터: 2019–2024 (1.47M행)  ·  예측 대상: 2025  ·  Strict Temporal Validation", 0.68, 5.80, 7.3, 0.30, size=11, color=MUTED)
    add_rect(slide, 0.65, 6.35, 7.10, 0.55, PANEL, radius=True)
    add_text(slide, "Phase 3 오프라인 해커톤 참가 여부  |  아니요", 0.88, 6.48, 6.62, 0.27, size=12, color=CORAL, bold=True)
    add_text(slide, "FINAL SOLUTION PRESENTATION", 9.05, 6.68, 3.33, 0.23, size=9, color="74809A", bold=True, align=PP_ALIGN.CENTER)

    # 02 — Problem Framing
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_header(slide, 2, "Problem framing", "확률을 맞히는 문제 — 순위보다 '잘 보정된 확률'")
    cards = [
        ("TARGET", "control_success", "각 투구의 제구 성공 확률\n0–1 실수 출력 (Brier 평가)", TEAL),
        ("METRIC", "Brier Skill Score", "1 - (Brier / Base_Brier)\n점수가 높을수록 우수", GOLD),
        ("UNIT", "1 pitch = 1 row", "평가 행별 100% 독립 예측\n다른 test 행 참조 절대 금지", BLUE),
    ]
    for i, (tag, value, note, color) in enumerate(cards):
        x = 0.58 + i * 4.15
        add_rect(slide, x, 2.25, 3.78, 2.10, PANEL, radius=True, line=GRID)
        add_text(slide, tag, x + 0.24, 2.50, 3.15, 0.25, size=10, color=color, bold=True)
        add_text(slide, value, x + 0.24, 2.88, 3.25, 0.42, size=20, bold=True)
        add_text(slide, note, x + 0.24, 3.45, 3.18, 0.60, size=12, color=MUTED)
    add_rect(slide, 0.58, 4.72, 12.15, 1.72, "10182C", radius=True, line=GRID)
    add_text(slide, "핵심 설계 원칙", 0.88, 4.98, 1.55, 0.30, size=13, color=CORAL, bold=True)
    principles = [
        "미래 시즌 예측 구조를 반영한 시간 순 엄격 검증",
        "공식 데이터만 사용 · 외부 API/데이터 원천 차단",
        "현재 행 + 사전 고정 모델만으로 완결되는 추론",
    ]
    for i, text_value in enumerate(principles):
        x = 2.35 + i * 3.32
        add_text(slide, f"0{i+1}", x, 4.98, 0.35, 0.30, size=11, color=TEAL, bold=True)
        add_text(slide, text_value, x + 0.45, 4.94, 2.62, 0.72, size=12, color=INK, bold=True)
    add_text(slide, "핵심 질문  |  '현재 투구 직전 정보만으로, 2025년 미래 시즌에서도 안정적이고 우수한 확률 품질을 달성할 수 있는가?'", 0.88, 6.00, 11.25, 0.25, size=11, color=MUTED)
    add_footer(slide, 2)

    # 03 — Data & Strict Temporal Validation
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_header(slide, 3, "Data & validation", "랜덤 K-Fold 분할을 지양하고 '과거 → 미래'로 엄격 검증")
    add_text(slide, "시계열 전이 검증 (Strict Temporal Forward Split)", 0.63, 2.22, 6.0, 0.30, size=12, color=TEAL, bold=True)
    y = 2.83
    x0 = 0.90
    seasons = list(range(2019, 2025))
    for i, season in enumerate(seasons):
        x = x0 + i * 1.73
        fill = TEAL_DARK if season <= 2023 else CORAL
        add_rect(slide, x, y, 1.28, 0.72, fill, radius=True)
        add_text(slide, str(season), x, y + 0.17, 1.28, 0.30, size=14, bold=True, align=PP_ALIGN.CENTER)
        if i < len(seasons)-1:
            add_line(slide, x + 1.28, y + 0.36, x + 1.70, y + 0.36, GRID, 1.4)
    add_text(slide, "TRAIN (2019–2023)  |  1,475,092 rows", 1.15, 3.76, 8.20, 0.28, size=12, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "VALID (2024)\n253,507 rows", 10.40, 3.70, 1.60, 0.55, size=11, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
    add_line(slide, 0.64, 4.35, 12.68, 4.35, GRID, 0.8)
    add_text(slide, "최종 학습 / 2025 실전 추론", 0.63, 4.70, 3.0, 0.30, size=12, color=GOLD, bold=True)
    for i, season in enumerate(seasons):
        x = x0 + i * 1.45
        add_rect(slide, x, 5.20, 1.05, 0.62, TEAL_DARK, radius=True)
        add_text(slide, str(season), x, 5.36, 1.05, 0.24, size=12, bold=True, align=PP_ALIGN.CENTER)
        if i < len(seasons)-1:
            add_line(slide, x + 1.05, 5.51, x + 1.42, 5.51, GRID, 1.3)
    add_text(slide, "→", 9.68, 5.26, 0.52, 0.36, size=22, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 10.25, 5.08, 1.95, 0.86, PANEL_2, radius=True, line=GOLD, width=1.5)
    add_text(slide, "2025 TEST\nPREDICT", 10.25, 5.22, 1.95, 0.52, size=13, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.66, 6.28, 11.97, 0.48, PANEL, radius=True)
    add_text(slide, "검증 원칙  |  다년도 야구 데이터에서 랜덤 K-Fold는 미래 누수를 유발하므로, T_train < T_val 시계열 전이 검증만 채택", 0.87, 6.40, 11.52, 0.24, size=11, color=MUTED)
    add_footer(slide, 3)

    # 04 — Feature Engineering
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_header(slide, 4, "Feature engineering", "공식 47개 + 파생 24개 + 집계 2개 + Trackman 8개 = 81개 핵심 피처")
    add_rect(slide, 0.60, 2.22, 3.15, 4.35, PANEL, radius=True, line=GRID)
    add_text(slide, "47", 0.86, 2.52, 1.10, 0.70, size=38, color=TEAL, bold=True)
    add_text(slide, "공식 입력 피처", 1.83, 2.72, 1.60, 0.32, size=14, bold=True)
    add_bullet_list(slide, ["경기·카운트·점수차", "주자·이닝·레버리지(li)", "투수/타자/팀/손잡이", "공식 asof_* 과거 이력"], 0.85, 3.54, 2.52, 2.20, size=13, gap=0.53)
    add_text(slide, "+", 3.86, 3.80, 0.50, 0.60, size=32, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 4.40, 2.22, 3.52, 4.35, PANEL, radius=True, line=GRID)
    add_text(slide, "24", 4.68, 2.52, 1.10, 0.70, size=38, color=GOLD, bold=True)
    add_text(slide, "상태·상호작용 파생", 5.72, 2.72, 1.85, 0.32, size=14, bold=True)
    feature_groups = [
        ("상황 압박", "count_state · scoring_pos · li"),
        ("매치업", "좌/우 platoon · batter pressure"),
        ("최근 추세", "직전 1·3경기 성공률 delta"),
        ("Cold-start", "결측 개수 · 결측 여부 · 이력 0 플래그"),
    ]
    for i, (tag, desc) in enumerate(feature_groups):
        yy = 3.48 + i * 0.69
        add_label(slide, tag, 4.69, yy, 0.90, color=GOLD, fill="222C47")
        add_text(slide, desc, 5.75, yy - 0.01, 1.85, 0.55, size=10, color=INK, bold=True)
    add_text(slide, "=", 8.06, 3.80, 0.50, 0.60, size=32, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 8.62, 2.22, 4.10, 4.35, "11222D", radius=True, line=TEAL_DARK, width=1.5)
    add_text(slide, "81", 8.98, 2.42, 1.78, 0.90, size=52, color=TEAL, bold=True)
    add_text(slide, "TOTAL FEATURES", 10.45, 2.80, 2.00, 0.30, size=11, color=MUTED, bold=True)
    add_line(slide, 8.98, 3.42, 12.33, 3.42, TEAL_DARK, 1.1)
    add_bullet_list(slide, ["row_id · target 제외", "결측 자체를 정보로 유지", "test 통계·빈도·순위 미사용", "학습·추론 동일 함수/순서 보장"], 9.00, 3.78, 3.10, 2.25, size=12, gap=0.52)
    add_text(slide, "모든 파생값은 해당 단일 행의 입력 컬럼만으로 계산", 9.00, 6.03, 3.30, 0.28, size=10, color=TEAL, bold=True)
    add_footer(slide, 4)

    # 05 — Architecture: 3-Tier Multi-Family GBDT Super Ensemble
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_header(slide, 5, "Model architecture", "3-Tier Multi-Family GBDT Super Ensemble & Adaptive Gate")
    add_rect(slide, 0.62, 2.25, 3.60, 4.40, PANEL, radius=True, line=GRID)
    add_card_title(slide, "Tier 1: Base & Specialists", 0.85, 2.50, 3.10, color=BLUE)
    add_bullet_list(slide, [
        "CatBoost 36-Seed Base (3-Channel)",
        "Futures 16-Model Regime Specialists",
        "3-Subtype Multi-Task Classifiers",
        "Psychological Latent Ridge Regressor"
    ], 0.85, 3.00, 3.15, 3.40, size=11, gap=0.80)

    add_rect(slide, 4.60, 2.25, 3.80, 4.40, PANEL, radius=True, line=GRID)
    add_card_title(slide, "Tier 2: Multi-Family & Blending", 4.85, 2.50, 3.30, color=GOLD)
    add_bullet_list(slide, [
        "LightGBM 1군 Regular R-Expert",
        "XGBoost Multi-Family Regressor",
        "Direct Brier Simplex Optimization",
        "Disjoint Matchup Empirical Bayes (113A)"
    ], 4.85, 3.00, 3.30, 3.40, size=11, gap=0.80)

    add_rect(slide, 8.80, 2.25, 3.90, 4.40, "11222D", radius=True, line=TEAL, width=1.5)
    add_card_title(slide, "Tier 3: Adaptive Meta Gate", 9.05, 2.50, 3.40, color=TEAL)
    add_bullet_list(slide, [
        "Zero-Centered Adaptive Meta Gate (scale=0.05)",
        "카운트/레버리지/상황 적응형 비선형 게이팅",
        "1군/2군 Macro-Leap Decoupling",
        "Global Calibration Shift (+0.0052)"
    ], 9.05, 3.00, 3.40, 3.40, size=11, gap=0.80)
    add_footer(slide, 5)

    # 06 — Decoupling & Disjoint EB Engine
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_header(slide, 6, "Specialization engine", "1군·2군 완전 디커플링 및 Disjoint Matchup EB 엔진")
    add_rect(slide, 0.62, 2.22, 5.80, 4.45, PANEL, radius=True, line=GRID)
    add_card_title(slide, "1군/2군 Macro-Leap Decoupling", 0.88, 2.50, 5.20, color=TEAL)
    add_bullet_list(slide, [
        "1군(Regular ~88%)과 2군(Futures ~12%)의 리그 성격 완전 분리",
        "1군: 36 CatBoost Base + LightGBM Expert + Adaptive Gate",
        "2군: 16 Futures CatBoost + 3 Subtypes + Psych Latent",
        "리그 간 혼선 방지로 2군 OOF BSS +3172.81pt 폭증 달성",
        "2025년 신인/이적 선수 분포 왜곡 사전 차단"
    ], 0.88, 3.10, 5.20, 3.30, size=11, gap=0.60)

    add_rect(slide, 6.82, 2.22, 5.85, 4.45, PANEL, radius=True, line=GRID)
    add_card_title(slide, "Disjoint Matchup Empirical Bayes (113A)", 7.08, 2.50, 5.20, color=GOLD)
    add_bullet_list(slide, [
        "투수-타자 상대전적 불균형에 대한 수축(Shrinkage) 추정",
        "소표본 매치업 과적합을 방지하는 정규화 EB 룩업",
        "Brier Simplex 가중치 최적화 결합 (w_eb = 0.035)",
        "단일 행 격리 추론 및 완전 무누수 사전 생성",
        "최종 리더보드 1121.90399점 달성의 핵심 기여 엔진"
    ], 7.08, 3.10, 5.20, 3.30, size=11, gap=0.60)
    add_footer(slide, 6)

    # 07 — Leaderboard Progression
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_header(
        slide,
        7,
        "Leaderboard progression",
        "815점 베이스라인에서 1121.90점 최종 챔피언까지의 진화",
        "공식 최고 점수: 1121.9039933605점 (submit_ref4_super113A.zip / Public 180위)",
    )
    headers = ["STAGE / VERSION", "KEY MECHANISM", "LB SCORE", "DELTA"]
    xs = [0.80, 4.20, 9.20, 11.20]
    ws = [3.20, 4.80, 1.80, 1.40]
    for t, x, w in zip(headers, xs, ws):
        add_text(slide, t, x, 2.20, w, 0.24, size=9, color=MUTED, bold=True)
    
    stages = [
        ("SUB-001 (Baseline)", "공식 47개 피처 LightGBM 베이스라인", "815.20127", "기준선", MUTED),
        ("SUB-002 (Regime R)", "Regime R-Capacity + Platt Calibration", "886.24881", "+71.0475", MUTED),
        ("EXP-030 (Champion Stack)", "3-Channel 6-Seed Residual + Shift (+0.0052)", "1068.25021", "+182.0014", TEAL),
        ("EXP-071 (Adaptive Gate)", "Zero-Centered Adaptive Multi-Channel Gate", "1092.18790", "+23.9377", TEAL),
        ("EXP-102 (Deep Hierarchical)", "Deep 61 Features + Multi-Seed Specialist", "1105.82017", "+13.6323", BLUE),
        ("EXP-107 (Super Ensemble)", "Multi-Family Simplex Super Blend (Top 174)", "1115.25607", "+9.4359", GOLD),
        ("EXP-109C (Tri-Family)", "Hyper-Regime Tri-Bridge 15-Model Tri-Family", "1120.89145", "+5.6354", GOLD),
        ("EXP-113A (Final Champion)", "Disjoint Matchup Empirical Bayes + 112C Base", "1121.90399", "+1.0125", CORAL),
    ]
    for i, (name, mech, score, delta, color) in enumerate(stages):
        yy = 2.65 + i * 0.52
        if i == len(stages) - 1:
            add_rect(slide, 0.65, yy - 0.05, 12.00, 0.44, "1F2E3B", radius=True, line=CORAL)
        add_text(slide, name, xs[0], yy, ws[0], 0.24, size=10, color=color, bold=(i >= 5))
        add_text(slide, mech, xs[1], yy, ws[1], 0.24, size=9, color=INK)
        add_text(slide, score, xs[2], yy, ws[2], 0.24, size=10, color=color, bold=True)
        add_text(slide, delta, xs[3], yy, ws[3], 0.24, size=9, color=color, bold=True)
    add_footer(slide, 7)

    # 08 — Inference Integrity & Row-Independence
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_header(slide, 8, "Inference integrity", "평가 데이터의 각 행을 100% 완전 독립적으로 예측")
    for i in range(4):
        y = 2.55 + i * 0.70
        add_rect(slide, 0.68, y, 2.42, 0.48, PANEL_2, radius=True, line=GRID)
        add_text(slide, f"row {i+1:02d}  ·  current inputs", 0.86, y + 0.12, 1.98, 0.20, size=9, color=INK, font=FONT_MONO)
        add_line(slide, 3.11, y + 0.24, 3.67, y + 0.24, TEAL_DARK, 1.2)
        add_rect(slide, 3.67, y, 2.30, 0.48, "11292C", radius=True)
        add_text(slide, "build_features(row)", 3.78, y + 0.12, 2.10, 0.20, size=9, color=TEAL, bold=True, font=FONT_MONO)
        add_line(slide, 5.98, y + 0.24, 6.42, y + 0.24, TEAL_DARK, 1.2)
        add_rect(slide, 6.42, y, 2.15, 0.48, PANEL, radius=True)
        add_text(slide, "fixed ensemble", 6.42, y + 0.12, 2.15, 0.20, size=9, color=INK, bold=True, align=PP_ALIGN.CENTER)
        add_line(slide, 8.58, y + 0.24, 9.12, y + 0.24, TEAL_DARK, 1.2)
        add_rect(slide, 9.12, y, 1.47, 0.48, "11222D", radius=True)
        add_text(slide, f"p{i+1}", 9.12, y + 0.12, 1.47, 0.20, size=10, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 10.85, 2.40, 1.85, 3.30, "241A25", radius=True, line=CORAL, width=1.2)
    add_text(slide, "금지 항목", 10.85, 2.65, 1.85, 0.35, size=16, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
    add_bullet_list(slide, ["test groupby", "rolling / shift", "batch statistics", "inference fit", "distribution scale"], 10.98, 3.20, 1.60, 2.30, size=9, bullet_color=CORAL, gap=0.38)
    add_rect(slide, 0.67, 5.75, 12.00, 0.75, PANEL, radius=True)
    checks = [
        ("배치 추론 = 단독 추론", "max diff: 0.000e+00 (PASS)"),
        ("행 순서 변경 (Shuffled)", "max diff: 0.000e+00 (PASS)"),
        ("단일 행 격리 추론", "max diff: 1.110e-16 (PASS)")
    ]
    for i, (title, note) in enumerate(checks):
        x = 0.92 + i * 4.00
        add_text(slide, "✓", x, 5.96, 0.32, 0.24, size=14, color=TEAL, bold=True)
        add_text(slide, title, x + 0.34, 5.92, 1.85, 0.26, size=10, color=INK, bold=True)
        add_text(slide, note, x + 0.34, 6.18, 2.50, 0.22, size=9, color=MUTED)
    add_footer(slide, 8)

    # 09 — Reproducibility & Fail-Fast Contract
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_header(slide, 9, "Reproducibility", "실행 환경·의존성·Fail-Fast 무결성 계약을 제출물에 고정")
    metrics = [
        ("1,475,092 rows", "전체 학습 행 수", TEAL),
        ("81 features", "최종 피처 계약", GOLD),
        ("Multi-Family", "CatBoost + LGBM + XGB", BLUE),
        ("AUDIT VERIFIED", "내부 무결성 감사", CORAL),
    ]
    for i, (value, label, color) in enumerate(metrics):
        x = 0.62 + i * 3.04
        add_rect(slide, x, 2.22, 2.70, 1.18, PANEL, radius=True, line=GRID)
        add_text(slide, value, x + 0.20, 2.52, 2.28, 0.34, size=16, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, label, x + 0.20, 2.96, 2.28, 0.22, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.62, 3.78, 5.92, 2.62, "10182C", radius=True, line=GRID)
    add_card_title(slide, "제출 실행 환경 명세", 0.92, 4.06, 2.50)
    env_rows = [
        ("Environment", "Ubuntu 22.04 / Python 3.11"),
        ("Engines", "CatBoost 1.2.10, LightGBM 4.0+, XGBoost 1.7+"),
        ("Requirements", "requirements.txt 완전 번들링"),
        ("Execution Time", "245K행 기준 ~1분 내외 완료 (10분 한도 여유)"),
    ]
    for i, (key, value) in enumerate(env_rows):
        yy = 4.53 + i * 0.42
        add_text(slide, key, 0.95, yy, 1.40, 0.22, size=8, color=MUTED, font=FONT_MONO)
        add_text(slide, value, 2.40, yy, 3.83, 0.22, size=9, color=INK, bold=True)
    add_rect(slide, 6.82, 3.78, 5.90, 2.62, "10182C", radius=True, line=GRID)
    add_card_title(slide, "Fail-Fast 제출 무결성 계약", 7.12, 4.06, 3.00, color=GOLD)
    add_bullet_list(slide, [
        "입력/피처 컬럼 개수, 순서, ID 계약 자동 검증",
        "NaN·Inf·범위 이탈([0, 1]) 예측 발생 시 즉시 실패",
        "출력 2개 컬럼(row_id, control_success) 및 입력 순서 100% 보존",
        "오프라인 단독 구동 및 네트워크 외부 호출 0건",
    ], 7.11, 4.51, 5.20, 1.75, size=10, bullet_color=GOLD, gap=0.41)
    add_rect(slide, 0.64, 6.64, 12.03, 0.31, "12262B", radius=True)
    add_text(slide, "ZIP 감사: AUDIT_VERIFIED  ·  행 독립성/재현성 검증 통과  ·  최종 제출: submit_ref4_super113A.zip", 0.87, 6.69, 11.60, 0.19, size=9, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 9)

    # 10 — Conclusion / Final Summary
    slide = prs.slides.add_slide(blank); set_bg(slide)
    add_text(slide, "10", 0.58, 0.45, 0.45, 0.28, size=11, color=TEAL, bold=True)
    add_text(slide, "CONCLUSION", 1.07, 0.45, 2.0, 0.28, size=10, color=MUTED, bold=True)
    add_text(slide, "미래 시즌에도 흔들리지 않는 제구 성공 확률 예측", 0.60, 1.03, 9.00, 0.62, size=28, bold=True)
    add_text(slide, "엄격한 시계열 검증, 완전한 행 독립성, 3-Tier Multi-Family GBDT로 완성", 0.62, 1.72, 9.00, 0.40, size=14, color=MUTED)
    takeaways = [
        ("01", "TIME-AWARE", "2019–2023 → 2024\n엄격한 시간 전이 검증"),
        ("02", "ROW-INDEPENDENT", "단일 행 독립 피처링\n오차 0.000e+00 달성"),
        ("03", "MULTI-FAMILY", "CatBoost + LGBM + XGB\nDisjoint EB Super Blend"),
        ("04", "FINAL 1121.90pt", "BSS 1121.903993점 달성\nPublic 180위 완주"),
    ]
    for i, (num, title, note) in enumerate(takeaways):
        x = 0.62 + i * 3.03
        add_rect(slide, x, 2.55, 2.70, 1.58, PANEL, radius=True, line=GRID)
        add_text(slide, num, x + 0.20, 2.78, 0.42, 0.28, size=11, color=TEAL, bold=True)
        add_text(slide, title, x + 0.63, 2.77, 1.80, 0.28, size=10, color=INK, bold=True)
        add_text(slide, note, x + 0.20, 3.30, 2.28, 0.54, size=11, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.62, 4.63, 12.10, 1.70, "241A25", radius=True, line=CORAL, width=1.6)
    add_text(slide, "Phase 3 오프라인 해커톤 참가 여부", 0.96, 4.99, 7.15, 0.36, size=17, color=INK, bold=True)
    add_text(slide, "아니요", 9.08, 4.87, 2.76, 0.66, size=34, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "본 발표자료는 온라인 해커톤(Phase 2) 최종 솔루션 설명 및 코드 검증용 제출 자료입니다.", 0.96, 5.62, 8.25, 0.28, size=11, color=MUTED)
    add_text(slide, "THANK YOU", 0.63, 6.72, 2.2, 0.27, size=10, color=TEAL, bold=True)
    add_text(slide, "LG AIMERS 9기 PHASE 2   |   최종 점수  1121.9039933605점 (180위)", 6.00, 6.69, 6.46, 0.28, size=10, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)
    add_footer(slide, 10)

    for idx, s in enumerate(prs.slides, start=1):
        s.name = f"Solution {idx:02d}"

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")
    
    SOLUTION_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(SOLUTION_PPTX)
    print(f"Saved: {SOLUTION_PPTX}")


if __name__ == "__main__":
    build_deck()

